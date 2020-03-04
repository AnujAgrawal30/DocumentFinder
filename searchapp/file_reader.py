from docx import Document
import io
from pdfminer.converter import TextConverter
from pdfminer.pdfinterp import PDFPageInterpreter
from pdfminer.pdfinterp import PDFResourceManager
from pdfminer.pdfpage import PDFPage
from pptx import Presentation


def read_docx(full_path):
    doc = Document(full_path)
    for para in doc.paragraphs:
        return para.text.strip().split(' ')


def read_txt(full_path):
    with open(full_path, 'rb') as f:
        try:
            text = f.read().decode().replace('\r', ' ').replace('\n', '')
            return text
        except UnicodeDecodeError:
            print("\rUnable to decode file at: {}".format(full_path))
        return ""


def read_pdf(pdf_path):
    resource_manager = PDFResourceManager()
    fake_file_handle = io.StringIO()
    converter = TextConverter(resource_manager, fake_file_handle)
    page_interpreter = PDFPageInterpreter(resource_manager, converter)

    with open(pdf_path, 'rb') as fh:
        for page in PDFPage.get_pages(fh,
                                      caching=True,
                                      check_extractable=True):
            page_interpreter.process_page(page)

        text = fake_file_handle.getvalue()
        # annots = pdfannots.process_file(fh, False)
        # if annots:
        #     if annots[0]:
        #         annotations = annots[0][0].gettext()

    # close open handles
    converter.close()
    fake_file_handle.close()

    if text:
        return text
        # return "Text:\n {} \n Annotations:\n {}".format(text, annotations)


def read_pptx(filepath):
    prs = Presentation(filepath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return text_runs
